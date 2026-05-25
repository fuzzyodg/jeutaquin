import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation(content_file, output_file):
    with open(content_file, 'r', encoding='utf-8') as f:
        slides_data = json.load(f)

    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Présentation Technique : Restaurant A'kadi"
    subtitle.text = "Architecture Tradi-Moderne : HTML, CSS & JavaScript\nAnalyse des Sections et du Code"

    for data in slides_data:
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Slide Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        title_tf = title_shape.text_frame
        title_tf.text = data['title']
        title_tf.paragraphs[0].font.size = Pt(28)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = RGBColor(111, 36, 10) # --color-primary

        # Image (Left)
        if os.path.exists(data.get('image', '')):
            slide.shapes.add_picture(data['image'], Inches(0.5), Inches(1), width=Inches(4.5))

        # Code (Right)
        code_text = ""
        if 'code_html' in data:
            code_text += "HTML:\n" + data['code_html'] + "\n\n"
        if 'code_css' in data:
            code_text += "CSS:\n" + data['code_css'] + "\n\n"
        if 'code_js' in data:
            code_text += "JS:\n" + data['code_js']

        code_shape = slide.shapes.add_textbox(Inches(5.2), Inches(1), Inches(4.3), Inches(4.5))
        code_tf = code_shape.text_frame
        code_tf.word_wrap = True
        p = code_tf.paragraphs[0]
        p.text = code_text
        p.font.size = Pt(10)
        p.font.name = 'Courier New'

        # Background for code box (light grey)
        fill = code_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 242, 235) # --color-light

        # Explanation (Bottom)
        expl_shape = slide.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(9), Inches(1.2))
        expl_tf = expl_shape.text_frame
        expl_tf.word_wrap = True
        p = expl_tf.paragraphs[0]
        p.text = data['explanation']
        p.font.size = Pt(14)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

if __name__ == "__main__":
    create_presentation('slides_content.json', 'presentation_akadi.pptx')
